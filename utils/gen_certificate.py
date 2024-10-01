"""
Module for generating and sending certificates via email.
"""

import os
import requests
from spire.presentation import Presentation, FileFormat
from pptx import Presentation as XPresentation
# from pptx.util import Inches
from config import Template
from utils.mail import send_email

PPT_DIR = "./files/ppt"
PDF_DIR = "./files/pdf"
TEMPLATES_DIR = "./files/templates"

for directory in [PPT_DIR, PDF_DIR, TEMPLATES_DIR]:
    os.makedirs(directory, exist_ok=True)

def delete_files(value):
    """
    Deletes files associated with a given value dictionary.
    
    Args:
        value (dict): A dictionary containing an 'id' key.
    """
    for extension in ['pptx', 'pdf']:
        try:
            os.remove(f'{PPT_DIR}/{value["id"]}.{extension}')
        except FileNotFoundError:
            pass  # Ignore if the file does not exist

def replace_text_in_presentation(presentation, value):
    """
    Replace placeholder text in the presentation with actual values.
    
    Args:
        presentation (XPresentation): The presentation to modify.
        value (dict): Values to replace in the presentation.
    """
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("{name}", value['name'])
                        run.text = run.text.replace("{id}", value['id'])

def generate_certificate(template: int, value, email: str = None):
    """
    Generates a certificate based on a given template and value,
    optionally sending it via email.
    
    Args:
        template (int): The index of the template to use.
        value (dict): Values to replace in the template.
        email (str, optional): Email to send the certificate to.
    
    Returns:
        dict: Result of the certificate generation and email sending process.
    """
    presentation = XPresentation(f'./templates/{Template[template]}.pptx')
    replace_text_in_presentation(presentation, value)
    presentation.save(f'{PPT_DIR}/{value["id"]}.pptx')
    result = ppt2pdf(value)
    delete_files(value)
    if email is None:
        return result
    try:
        with open("./templates/mail.html", "r", encoding="utf-8") as file:
            html_body = file.read().replace("${name}", value["name"]) \
                                     .replace("${body}", value["body"]) \
                                     .replace("${download}", result["dl"])
            res = send_email(email, value["subject"], html_body)
            res['dl'] = result['dl']
            return res
    except FileNotFoundError:
        delete_files(value)
        return {"success": False, "message": "Failed to read HTML file"}

def generate_certificate_from_url(template_url: str, value, email: str = None):
    """
    Generates a certificate from a given template URL and customizes it.
        
    Args:
        template_url (str): The URL of the certificate template.
        value (dict): Values to replace in the template.
        email (str, optional): Email to send the certificate to.
        
    Returns:
        dict: Result of the certificate generation and email sending process.
    """
    try:
        rs = requests.get(template_url, timeout=100)
        with open(f"{TEMPLATES_DIR}/{value['id']}.pptx", "wb") as file:
            file.write(rs.content)
        presentation = XPresentation(f"{TEMPLATES_DIR}/{value['id']}.pptx")
        replace_text_in_presentation(presentation, value)
        presentation.save(f"{PPT_DIR}/{value['id']}.pptx")
        result = ppt2pdf(value)
        delete_files(value)
        if email is None:
            return result
        with open("./templates/mail.html", "r", encoding="utf-8") as file:
            html_body = file.read().replace("${name}", value["name"]) \
                                     .replace("${body}", value["body"]) \
                                     .replace("${download}", result["dl"])
            res = send_email(email, value["subject"], html_body)
            res['dl'] = result['dl']
            return res
    except requests.RequestException as e:
        delete_files(value)
        return {"success": False, "message": f"Failed to fetch template: {e}"}

def ppt2pdf(value):
    """
    Converts a PowerPoint presentation to a PDF and uploads it to a remote server.
    
    Args:
        value (dict): A dictionary containing the 'id' of the PowerPoint file.
    
    Returns:
        dict: Success status and download link or error message.
    """
    presentation = Presentation()
    presentation.LoadFromFile(f"{PPT_DIR}/{value['id']}.pptx")
    presentation.Slides[0].SaveToFile(f"{PDF_DIR}/{value['id']}.pdf", FileFormat.PDF)
    presentation.Dispose()
    try:
        with open(f"{PDF_DIR}/{value['id']}.pdf", 'rb') as file:
            rs = requests.post("https://certify.izaries.workers.dev/certificate",
                               data=value,
                               files={'file': file}, timeout=100).json()
        if rs["success"]:
            return {
                "success": True,
                "dl": f"https://certify.izaries.workers.dev/download?id={rs['id']}"
            }
        return {"success": False, "message": rs["message"]}
    except requests.RequestException as e:
        return {"success": False, "message": f"Failed to upload PDF: {e}"}

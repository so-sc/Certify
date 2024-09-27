

import requests, os, io
from spire.presentation import *
from spire.presentation.common import *
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

from config import Template
from utils.mail import sendMail

if not os.path.exists("./files"):
    os.mkdir("./files")
ppt = "./files/ppt"
if not os.path.exists(ppt):
    os.mkdir(ppt)
pdf = "./files/pdf"
if not os.path.exists(pdf):
    os.mkdir(pdf)


def genCertificate(template: int, value, email: str = None):
    presentation = Presentation(f'./templates/{Template[template]}.pptx')
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{name}" in run.text:
                            run.text = run.text.replace("{name}", value['name'])
                        if "{id}" in run.text:
                            run.text = run.text.replace("{id}", value['id'])
    presentation.save(f'{ppt}/{value["id"]}.pptx')
    result = ppt2pdf(value)
    try:
        os.remove(f'{ppt}/{value["id"]}.pptx')
    except Exception:
        pass
    try:
        os.remove(f'{pdf}/{value["id"]}.pdf')
    except Exception:
        pass
    if not email:
        return result
    try:
        with open("./templates/mail.html", "r") as file:
            html_body = file.read()
            html_body = html_body.replace("${name}", value["name"])
            html_body = html_body.replace("${body}", value["body"])
            html_body = html_body.replace("${download}", result["dl"])
            res = sendMail(value["name"], email, value["subject"], html_body)
            res['dl'] = result['dl']
            return res
    except Exception as e:
        print(f"Failed to read HTML file: {e}")
        return { "success": False, "message": "Failed to read HTML file" }

def ppt2pdf(value):
    presentation = Presentation()
    presentation.LoadFromFile(f"{ppt}/{value['id']}.pptx")
    slide = presentation.Slides[0]
    slide.SaveToFile(f"{pdf}/{value['id']}.pdf", FileFormat.PDF)
    presentation.Dispose()
    rs = requests.post("https://certify.izaries.workers.dev/certificate",
        data=value,
        files={'file': open(f"{pdf}/{value['id']}.pdf", 'rb')}
    ).json()
    if rs["success"]:
        return { "success": True, "dl": f"https://certify.izaries.workers.dev/download?id={rs['id']}" }
    else:
        return { "success": False, "message": rs["message"] }


def genCertificateFromUrl(template_url: str, value, email: str = None):
    response = requests.get(template_url)
    if response.status_code != 200:
        return { "success": False, "message": "Failed to download the template" }
    
    pptx_data = io.BytesIO(response.content)
    presentation = Presentation(pptx_data)

    # Update the PPTX content
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{name}" in run.text:
                            run.text = run.text.replace("{name}", value['name'])
                        if "{id}" in run.text:
                            run.text = run.text.replace("{id}", value['id'])

    # Convert PPTX to PDF
    pdf_buffer = io.BytesIO()
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    for slide in presentation.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Add each shape's text to the PDF
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pdf.set_font("Arial", size=12)
                        pdf.multi_cell(0, 10, run.text)

    # Save the PDF to the buffer
    pdf_buffer.seek(0)  # Ensure buffer is at the beginning
    pdf.output(pdf_buffer, 'S')  # Write to buffer as string instead of file

    # Send the PDF to the server
    pdf_buffer.seek(0)  # Reset buffer position before sending
    rs = requests.post(
        "https://certify.izaries.workers.dev/certificate",
        data=value,
        files={'file': ('certificate.pdf', pdf_buffer, 'application/pdf')}
    ).json()

    if rs["success"]:
        return { "success": True, "dl": f"https://certify.izaries.workers.dev/download?id={rs['id']}" }
    else:
        return { "success": False, "message": rs["message"] }
